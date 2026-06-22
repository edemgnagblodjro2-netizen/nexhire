"""Export de rapports professionnels avec logo org — PDF / Excel / Word / PowerPoint."""
from __future__ import annotations

import io
from datetime import datetime
from typing import Literal

from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

from auth import CurrentUser
from db import get_db, row
from rbac import require_min_role

router = APIRouter(prefix="/api/agent", tags=["reports"])

BRAND  = "#818CF8"
NAVY   = "#0f172a"
GRAY   = "#64748b"
LIGHT  = "#e2e8f0"
BG     = "#f8fafc"
COLORS = ["#818CF8","#6366f1","#4f46e5","#0ea5e9","#10b981","#f59e0b","#ef4444","#94a3b8"]


# ── Models ─────────────────────────────────────────────────────────────────

class ChartSpec(BaseModel):
    type:   Literal["pie", "bar", "line"] = "bar"
    title:  str
    labels: list[str]
    values: list[float]


class ExportRequest(BaseModel):
    question:     str = Field(..., min_length=1, max_length=2000)
    answer:       str = Field(..., min_length=1, max_length=20000)
    sources:      list[str] = []
    format:       Literal["pdf", "xlsx", "docx", "pptx"] = "pdf"
    title:        str | None = None
    conclusion:   str | None = None
    charts:       list[ChartSpec] = []
    date_from:    str | None = None   # YYYY-MM-DD
    date_to:      str | None = None   # YYYY-MM-DD
    period_label: str | None = None   # ex: "Juin 2026" ou "3 derniers mois"


@router.post("/export")
def export_report(
    payload: ExportRequest,
    user: CurrentUser = Depends(require_min_role("user")),
):
    title    = payload.title or _derive_title(payload.question)
    date_str = datetime.now().strftime("%d %B %Y à %H:%M")

    org_name = "Organisation"
    logo_url: str | None = None
    if user.organization_id:
        with get_db() as cur:
            cur.execute(
                "SELECT name, logo_url FROM organizations WHERE id = %s LIMIT 1",
                (user.organization_id,),
            )
            r = row(cur) or {}
        org_name = r.get("name") or "Organisation"
        logo_url = r.get("logo_url") or None

    logo_bytes  = _fetch_logo(logo_url)
    budget_dept = _fetch_budget_by_dept(user.organization_id or "", payload.date_from, payload.date_to)
    budget_cat  = _fetch_budget_by_category(user.organization_id or "", payload.date_from, payload.date_to)

    if payload.format == "pdf":
        return _export_pdf(title, date_str, payload, org_name, logo_bytes, budget_dept, budget_cat)
    elif payload.format == "xlsx":
        return _export_excel(title, date_str, payload, org_name, budget_dept, budget_cat)
    elif payload.format == "docx":
        return _export_word(title, date_str, payload, org_name, logo_bytes, budget_dept, budget_cat)
    elif payload.format == "pptx":
        return _export_pptx(title, date_str, payload, org_name, logo_bytes, budget_dept, budget_cat)


# ── Shared helpers ─────────────────────────────────────────────────────────

def _derive_title(q: str) -> str:
    q = q.strip()
    return (q[0].upper() + q[1:])[:80] if q else "Rapport"


def _filename(fmt: str) -> str:
    return f"nexhire-rapport-{datetime.now().strftime('%Y%m%d-%H%M')}.{fmt}"


def _default_conclusion(payload: ExportRequest, date_str: str) -> str:
    srcs = ", ".join(s.upper() for s in payload.sources) if payload.sources else "—"
    return (
        f"Ce rapport a été généré automatiquement par NexHire EIP "
        f"le {date_str} à partir des données extraites de : {srcs}.\n"
        "Les informations sont confidentielles et réservées aux personnes autorisées.\n"
        "Pour toute question, contactez votre administrateur NexHire."
    )


def _fetch_budget_by_dept(org_id: str, date_from: str | None, date_to: str | None) -> list[dict]:
    """Top 10 départements par budget alloué (allocated / actual / balance)."""
    conds = ["b.organization_id = %s"]
    params: list = [org_id]
    _apply_date_filter(conds, params, date_from, date_to)
    where = " AND ".join(conds)
    sql = f"""
        SELECT
            COALESCE(d.name, 'Sans département') AS name,
            SUM(b.allocated)  AS allocated,
            SUM(b.actual)     AS actual,
            MAX(b.currency)   AS currency
        FROM budget_entries b
        LEFT JOIN departments d ON b.department_id = d.id
        WHERE {where}
        GROUP BY d.name
        ORDER BY allocated DESC
        LIMIT 10
    """
    try:
        with get_db() as cur:
            cur.execute(sql, params)
            return rows(cur)
    except Exception:
        return []


def _fetch_budget_by_category(org_id: str, date_from: str | None, date_to: str | None) -> list[dict]:
    """Top 10 catégories par dépenses réelles (actual)."""
    conds = ["b.organization_id = %s"]
    params: list = [org_id]
    _apply_date_filter(conds, params, date_from, date_to)
    where = " AND ".join(conds)
    sql = f"""
        SELECT
            b.category        AS name,
            SUM(b.allocated)  AS allocated,
            SUM(b.actual)     AS actual,
            MAX(b.currency)   AS currency
        FROM budget_entries b
        WHERE {where}
        GROUP BY b.category
        ORDER BY actual DESC
        LIMIT 10
    """
    try:
        with get_db() as cur:
            cur.execute(sql, params)
            return rows(cur)
    except Exception:
        return []


def _apply_date_filter(conds: list, params: list, date_from: str | None, date_to: str | None) -> None:
    if date_from:
        try:
            y, m = int(date_from[:4]), int(date_from[5:7])
            conds.append("(b.year > %s OR (b.year = %s AND (b.month IS NULL OR b.month >= %s)))")
            params.extend([y, y, m])
        except (ValueError, IndexError):
            pass
    if date_to:
        try:
            y, m = int(date_to[:4]), int(date_to[5:7])
            conds.append("(b.year < %s OR (b.year = %s AND (b.month IS NULL OR b.month <= %s)))")
            params.extend([y, y, m])
        except (ValueError, IndexError):
            pass


def _fetch_logo(url: str | None) -> bytes | None:
    """Télécharge le logo (PNG/JPG seulement — pas SVG pour compatibilité PDF/PPTX)."""
    if not url:
        return None
    ext = url.lower().split("?")[0].rsplit(".", 1)[-1]
    if ext == "svg":
        return None
    try:
        import httpx
        r = httpx.get(url, timeout=5, follow_redirects=True)
        ct = r.headers.get("content-type", "")
        if r.status_code == 200 and ("image/png" in ct or "image/jpeg" in ct):
            return r.content
    except Exception:
        pass
    return None


def _chart_image(spec: ChartSpec, w: float = 6.5, h: float = 3.6) -> bytes:
    """Génère un graphique matplotlib en PNG (bytes)."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#f8faff")

    n = len(spec.labels)
    palette = COLORS[:n] if n <= len(COLORS) else COLORS * ((n // len(COLORS)) + 1)

    if spec.type == "pie":
        wedges, texts, autotexts = ax.pie(
            spec.values, labels=spec.labels, autopct="%1.1f%%",
            colors=palette[:n], startangle=90,
            pctdistance=0.78, wedgeprops={"linewidth": 1.5, "edgecolor": "white"},
        )
        for t in texts:     t.set_fontsize(9)
        for t in autotexts: t.set_fontsize(8); t.set_color("white"); t.set_fontweight("bold")

    elif spec.type in ("bar", "line"):
        x = range(n)
        if spec.type == "bar":
            bars = ax.bar(x, spec.values, color=palette[:n], width=0.55,
                          edgecolor="white", linewidth=1.2)
            for bar, val in zip(bars, spec.values):
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(spec.values) * 0.01,
                        f"{val:g}", ha="center", va="bottom", fontsize=8.5, fontweight="bold",
                        color="#0f172a")
        else:
            ax.plot(list(x), spec.values, color=COLORS[0], linewidth=2.5,
                    marker="o", markersize=6, markerfacecolor="white", markeredgewidth=2)
            ax.fill_between(list(x), spec.values, alpha=0.12, color=COLORS[0])

        ax.set_xticks(list(x))
        ax.set_xticklabels(spec.labels, fontsize=8.5, rotation=15 if n > 4 else 0, ha="right")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#e2e8f0")
        ax.spines["bottom"].set_color("#e2e8f0")
        ax.tick_params(colors="#64748b")
        ax.yaxis.set_tick_params(labelsize=8)
        ax.grid(axis="y", linestyle="--", alpha=0.4, color="#e2e8f0")

    ax.set_title(spec.title, fontsize=11, fontweight="bold", color="#0f172a", pad=10)
    plt.tight_layout(pad=0.8)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── PDF ────────────────────────────────────────────────────────────────────

def _export_pdf(
    title: str,
    date_str: str,
    payload: ExportRequest,
    org_name: str,
    logo_bytes: bytes | None,
    budget_dept: list[dict],
    budget_cat: list[dict],
) -> StreamingResponse:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.lib.units import cm
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import BaseDocTemplate, PageTemplate, Frame, Paragraph, Spacer, HRFlowable, Image as RLImage

    w_page, h_page = A4
    HEADER_H = 1.55 * cm
    FOOTER_H = 1.1  * cm
    MARGIN   = 2.0  * cm

    C_BRAND = HexColor(BRAND)
    C_NAVY  = HexColor(NAVY)
    C_GRAY  = HexColor(GRAY)
    C_LIGHT = HexColor(LIGHT)
    C_BG    = HexColor(BG)
    C_WHITE = HexColor("#ffffff")

    logo_reader: ImageReader | None = None
    if logo_bytes:
        try:
            logo_reader = ImageReader(io.BytesIO(logo_bytes))
        except Exception:
            logo_reader = None

    def _on_page(canvas, doc):
        canvas.saveState()

        # ── HEADER bar
        canvas.setFillColor(C_NAVY)
        canvas.rect(0, h_page - HEADER_H, w_page, HEADER_H, fill=1, stroke=0)

        # Accent stripe (brand colour, left edge)
        canvas.setFillColor(C_BRAND)
        canvas.rect(0, h_page - HEADER_H, 0.35 * cm, HEADER_H, fill=1, stroke=0)

        logo_end_x = 0.9 * cm
        if logo_reader:
            logo_h = HEADER_H - 0.5 * cm
            logo_w = logo_h * 4  # generous max width; preserveAspectRatio clips it
            canvas.drawImage(
                logo_reader,
                0.8 * cm, h_page - HEADER_H + 0.25 * cm,
                width=logo_w, height=logo_h,
                preserveAspectRatio=True, mask="auto",
            )
            logo_end_x = 0.8 * cm + logo_w + 0.3 * cm
        else:
            # Text fallback: org name
            canvas.setFont("Helvetica-Bold", 11)
            canvas.setFillColor(C_WHITE)
            nex = "Nex"
            nex_w = canvas.stringWidth(nex, "Helvetica-Bold", 11)
            canvas.drawString(0.85 * cm, h_page - HEADER_H + 0.52 * cm, nex)
            canvas.setFillColor(C_BRAND)
            canvas.drawString(0.85 * cm + nex_w, h_page - HEADER_H + 0.52 * cm, "hire EIP")

        # RAPPORT + org name (top-right)
        canvas.setFont("Helvetica-Bold", 7)
        canvas.setFillColor(C_WHITE)
        canvas.drawRightString(w_page - 0.7 * cm, h_page - 0.65 * cm, "RAPPORT")
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(C_BRAND)
        canvas.drawRightString(w_page - 0.7 * cm, h_page - 1.05 * cm, org_name[:50])

        # ── FOOTER bar
        canvas.setFillColor(C_BG)
        canvas.rect(0, 0, w_page, FOOTER_H, fill=1, stroke=0)
        canvas.setStrokeColor(C_LIGHT)
        canvas.line(0, FOOTER_H, w_page, FOOTER_H)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(C_GRAY)
        footer_left = f"{org_name} · NexHire EIP · Confidentiel — usage interne uniquement"
        if payload.period_label:
            footer_left = f"{org_name} · {payload.period_label} · Confidentiel"
        canvas.drawString(0.7 * cm, 0.38 * cm, footer_left)
        canvas.drawRightString(w_page - 0.7 * cm, 0.38 * cm, f"Page {doc.page}")

        canvas.restoreState()

    buf = io.BytesIO()
    content_y      = FOOTER_H + 0.4 * cm
    content_height = h_page - HEADER_H - FOOTER_H - 0.8 * cm
    frame = Frame(MARGIN, content_y, w_page - 2 * MARGIN, content_height, id="main")
    tmpl  = PageTemplate(id="main", frames=[frame], onPage=_on_page)
    doc   = BaseDocTemplate(buf, pagesize=A4, pageTemplates=[tmpl])

    s_title  = ParagraphStyle("title",  fontName="Helvetica-Bold", fontSize=18, textColor=C_NAVY, spaceAfter=4, spaceBefore=6)
    s_date   = ParagraphStyle("date",   fontName="Helvetica",       fontSize=10, textColor=C_GRAY, spaceAfter=10)
    s_label  = ParagraphStyle("label",  fontName="Helvetica-Bold",  fontSize=9,  textColor=C_BRAND, spaceBefore=14, spaceAfter=4)
    s_body   = ParagraphStyle("body",   fontName="Helvetica",       fontSize=11, textColor=C_NAVY, leading=18, spaceAfter=4)
    s_concl  = ParagraphStyle("concl",  fontName="Helvetica",       fontSize=10, textColor=C_GRAY, leading=15, spaceAfter=4)
    s_src    = ParagraphStyle("src",    fontName="Helvetica",       fontSize=10, textColor=C_GRAY, spaceAfter=6)

    conclusion = payload.conclusion or _default_conclusion(payload, date_str)
    story: list = []

    # ── Title block
    story.append(Paragraph(title, s_title))
    period_line = f"Période : {payload.period_label}   ·   Généré le {date_str}" if payload.period_label else f"Généré le {date_str}"
    story.append(Paragraph(period_line, s_date))
    if payload.sources:
        story.append(Paragraph("SOURCES CONSULTÉES", s_label))
        story.append(Paragraph(" &nbsp;·&nbsp; ".join(s.upper() for s in payload.sources), s_src))
    story.append(HRFlowable(color=C_LIGHT, thickness=1, spaceAfter=8))

    # ── Body
    story.append(Paragraph("QUESTION", s_label))
    story.append(Paragraph(payload.question.replace("\n", "<br/>"), s_body))
    story.append(Paragraph("RÉPONSE DE L'AGENT", s_label))
    story.append(HRFlowable(color=C_LIGHT, thickness=1, spaceAfter=6))
    for line in payload.answer.split("\n"):
        if line.strip():
            story.append(Paragraph(line.replace("&", "&amp;").replace("<", "&lt;"), s_body))

    # ── Charts
    if payload.charts:
        story.append(Paragraph("GRAPHIQUES", s_label))
        for spec in payload.charts:
            try:
                img = RLImage(io.BytesIO(_chart_image(spec)), width=14 * cm, height=7.8 * cm)
                story.append(img)
                story.append(Spacer(1, 0.3 * cm))
            except Exception:
                pass

    # ── Budget par département — Top 10
    if budget_dept:
        from reportlab.platypus import Table, TableStyle
        currency = budget_dept[0].get("currency") or "CAD"
        total_alloc  = sum(r.get("allocated") or 0 for r in budget_dept)
        total_actual = sum(r.get("actual")    or 0 for r in budget_dept)

        story.append(Spacer(1, 0.5 * cm))
        story.append(HRFlowable(color=C_LIGHT, thickness=1, spaceAfter=6))
        story.append(Paragraph("BUDGET PAR DÉPARTEMENT — TOP 10", s_label))
        story.append(Paragraph(
            f"Total alloué : <b>{total_alloc:,.0f} {currency}</b>  ·  "
            f"Total réel : <b>{total_actual:,.0f} {currency}</b>  ·  "
            f"Écart : <b>{total_alloc - total_actual:,.0f} {currency}</b>",
            ParagraphStyle("bsum", fontName="Helvetica", fontSize=10, textColor=C_NAVY, spaceAfter=8, leading=15),
        ))

        tdata = [["Département", "Alloué", "Réel / Dépenses", "Écart", "% Consommé"]]
        bar_labels, bar_pct = [], []
        for r in budget_dept:
            alloc  = r.get("allocated") or 0
            actual = r.get("actual")    or 0
            bal    = alloc - actual
            pct    = actual / alloc * 100 if alloc > 0 else 0
            dept_name = (r.get("name") or "—")[:28]
            tdata.append([dept_name, f"{alloc:,.0f}", f"{actual:,.0f}", f"{bal:,.0f}", f"{pct:.1f}%"])
            bar_labels.append(dept_name); bar_pct.append(round(pct, 1))

        tbl = Table(tdata, colWidths=[5.8*cm, 3*cm, 3.2*cm, 3*cm, 2.5*cm])
        tbl.setStyle(TableStyle([
            ("BACKGROUND",    (0, 0), (-1, 0),  C_NAVY),
            ("TEXTCOLOR",     (0, 0), (-1, 0),  HexColor("#ffffff")),
            ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
            ("FONTSIZE",      (0, 0), (-1, -1), 8.5),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1), [HexColor("#f0f4ff"), HexColor("#ffffff")]),
            ("GRID",          (0, 0), (-1, -1), 0.4, C_LIGHT),
            ("ALIGN",         (1, 0), (-1, -1), "RIGHT"),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING",    (0, 0), (-1, -1), 4),
        ]))
        story.append(tbl)
        story.append(Spacer(1, 0.3 * cm))

        try:
            spec_dept = ChartSpec(type="bar", title="% Budget consommé par département", labels=bar_labels, values=bar_pct)
            story.append(RLImage(io.BytesIO(_chart_image(spec_dept, w=6.5, h=3.2)), width=13*cm, height=6.4*cm))
        except Exception:
            pass

    # ── Budget par catégorie — Top 10
    if budget_cat:
        currency = budget_cat[0].get("currency") or "CAD"
        story.append(Spacer(1, 0.4 * cm))
        story.append(Paragraph("DÉPENSES PAR CATÉGORIE — TOP 10", s_label))

        tdata2 = [["Catégorie", "Alloué", "Réel / Dépenses", "% du total réel"]]
        total_actual_cat = sum(r.get("actual") or 0 for r in budget_cat) or 1
        cat_labels, cat_vals = [], []
        for r in budget_cat:
            alloc  = r.get("allocated") or 0
            actual = r.get("actual")    or 0
            pct    = actual / total_actual_cat * 100
            cat_name = (r.get("name") or "—")[:28]
            tdata2.append([cat_name, f"{alloc:,.0f}", f"{actual:,.0f}", f"{pct:.1f}%"])
            cat_labels.append(cat_name); cat_vals.append(round(actual, 0))

        tbl2 = Table(tdata2, colWidths=[5.8*cm, 3*cm, 3.2*cm, 3*cm])
        tbl2.setStyle(TableStyle([
            ("BACKGROUND",    (0, 0), (-1, 0),  C_NAVY),
            ("TEXTCOLOR",     (0, 0), (-1, 0),  HexColor("#ffffff")),
            ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
            ("FONTSIZE",      (0, 0), (-1, -1), 8.5),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1), [HexColor("#f0f4ff"), HexColor("#ffffff")]),
            ("GRID",          (0, 0), (-1, -1), 0.4, C_LIGHT),
            ("ALIGN",         (1, 0), (-1, -1), "RIGHT"),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING",    (0, 0), (-1, -1), 4),
        ]))
        story.append(tbl2)
        story.append(Spacer(1, 0.3 * cm))

        try:
            spec_cat = ChartSpec(type="pie", title="Répartition des dépenses par catégorie", labels=cat_labels, values=cat_vals)
            story.append(RLImage(io.BytesIO(_chart_image(spec_cat, w=5.5, h=3.5)), width=11*cm, height=7*cm))
        except Exception:
            pass

    # ── Conclusion
    story.append(Spacer(1, 0.4 * cm))
    story.append(HRFlowable(color=C_BRAND, thickness=1.5, spaceAfter=6))
    story.append(Paragraph("CONCLUSION", s_label))
    for line in conclusion.split("\n"):
        if line.strip():
            story.append(Paragraph(line.replace("&", "&amp;").replace("<", "&lt;"), s_concl))

    doc.build(story)
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/pdf",
                             headers={"Content-Disposition": f'attachment; filename="{_filename("pdf")}"'})


# ── Excel ──────────────────────────────────────────────────────────────────

def _export_excel(
    title: str,
    date_str: str,
    payload: ExportRequest,
    org_name: str,
    budget_dept: list[dict],
    budget_cat: list[dict],
) -> StreamingResponse:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.chart import BarChart, PieChart, LineChart, Reference

    wb = Workbook()
    ws = wb.active
    ws.title = "Rapport"
    ws.column_dimensions["A"].width = 24
    ws.column_dimensions["B"].width = 88

    NAVY_XL  = "0F172A"; BRAND_XL = "818CF8"; GRAY_XL = "64748B"; SLATE_XL = "475569"

    def hdr(cell, text, bg, fg="FFFFFF", sz=11, bold=True):
        cell.value = text
        cell.font = Font(bold=bold, size=sz, color=fg)
        cell.fill = PatternFill(fill_type="solid", fgColor=bg)
        cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def body(cell, text):
        cell.value = text
        cell.alignment = Alignment(wrap_text=True, vertical="top")
        cell.font = Font(size=10, color=NAVY_XL)

    # EN-TÊTE
    ws.merge_cells("A1:B1")
    ws["A1"].value = f"{org_name}  ·  NexHire EIP"
    ws["A1"].font  = Font(bold=True, size=15, color="FFFFFF")
    ws["A1"].fill  = PatternFill(fill_type="solid", fgColor=NAVY_XL)
    ws["A1"].alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 36

    ws.merge_cells("A2:B2")
    ws["A2"].value = title
    ws["A2"].font  = Font(bold=True, size=13, color=NAVY_XL)
    ws["A2"].alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[2].height = 26

    ws.merge_cells("A3:B3")
    period_cell = f"Période : {payload.period_label}   ·   Généré le {date_str}" if payload.period_label else f"Généré le {date_str}"
    ws["A3"].value = period_cell
    ws["A3"].font  = Font(size=9, italic=True, color=GRAY_XL)
    ws.row_dimensions[3].height = 18

    row = 5
    if payload.sources:
        hdr(ws[f"A{row}"], "Sources", BRAND_XL)
        body(ws[f"B{row}"], " · ".join(s.upper() for s in payload.sources))
        ws.row_dimensions[row].height = 22; row += 1

    hdr(ws[f"A{row}"], "Question", SLATE_XL)
    body(ws[f"B{row}"], payload.question)
    ws.row_dimensions[row].height = max(30, 15 * (len(payload.question) // 80 + 1)); row += 1

    answer_lines = [l for l in payload.answer.split("\n") if l.strip()]
    hdr(ws[f"A{row}"], "Réponse de l'agent", SLATE_XL)
    ws[f"A{row}"].alignment = Alignment(wrap_text=True, vertical="top")
    body(ws[f"B{row}"], "\n".join(answer_lines))
    ws.row_dimensions[row].height = max(60, min(400, 14 * len(answer_lines))); row += 2

    for spec in payload.charts:
        chart_ws = wb.create_sheet(title=spec.title[:30])
        chart_ws.column_dimensions["A"].width = 28
        chart_ws.column_dimensions["B"].width = 16
        chart_ws["A1"] = "Catégorie"; chart_ws["B1"] = "Valeur"
        chart_ws["A1"].font = Font(bold=True, color="FFFFFF")
        chart_ws["A1"].fill = PatternFill(fill_type="solid", fgColor=NAVY_XL)
        chart_ws["B1"].font = chart_ws["A1"].font
        chart_ws["B1"].fill = chart_ws["A1"].fill

        for i, (lbl, val) in enumerate(zip(spec.labels, spec.values), start=2):
            chart_ws[f"A{i}"] = lbl; chart_ws[f"B{i}"] = val

        n = len(spec.labels)
        data_ref = Reference(chart_ws, min_col=2, min_row=1, max_row=n + 1)
        cats_ref = Reference(chart_ws, min_col=1, min_row=2, max_row=n + 1)
        chart = PieChart() if spec.type == "pie" else (LineChart() if spec.type == "line" else BarChart())
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)
        chart.title = spec.title
        chart.shape = 4
        chart_ws.add_chart(chart, "D2")

    conclusion = payload.conclusion or _default_conclusion(payload, date_str)
    hdr(ws[f"A{row}"], "Conclusion", BRAND_XL)
    ws[f"A{row}"].alignment = Alignment(wrap_text=True, vertical="top")
    body(ws[f"B{row}"], conclusion)
    ws.row_dimensions[row].height = max(50, min(200, 14 * (conclusion.count("\n") + 3))); row += 2

    ws.merge_cells(f"A{row}:B{row}")
    ws[f"A{row}"].value = f"{org_name} · NexHire EIP · Confidentiel · {date_str}"
    ws[f"A{row}"].font  = Font(size=8, italic=True, color=GRAY_XL)

    # ── Feuille Budget Département
    if budget_dept:
        currency = budget_dept[0].get("currency") or "CAD"
        bws = wb.create_sheet(title="Budget Département")
        bws.column_dimensions["A"].width = 30
        for col in "BCDE": bws.column_dimensions[col].width = 18

        def _bhdr(cell, text, bg=NAVY_XL, fg="FFFFFF"):
            cell.value = text
            cell.font  = Font(bold=True, size=10, color=fg)
            cell.fill  = PatternFill(fill_type="solid", fgColor=bg)
            cell.alignment = Alignment(horizontal="center", vertical="center")

        bws.merge_cells("A1:E1")
        bws["A1"].value = f"Budget par département — Top 10   ·   {org_name}"
        bws["A1"].font  = Font(bold=True, size=13, color="FFFFFF")
        bws["A1"].fill  = PatternFill(fill_type="solid", fgColor=NAVY_XL)
        bws["A1"].alignment = Alignment(horizontal="left", vertical="center")
        bws.row_dimensions[1].height = 28

        for col, label in zip("ABCDE", ["Département", f"Alloué ({currency})", f"Réel ({currency})", f"Écart ({currency})", "% Consommé"]):
            _bhdr(bws[f"{col}2"], label, BRAND_XL)
        bws.row_dimensions[2].height = 22

        fills = [PatternFill(fill_type="solid", fgColor="F0F4FF"), PatternFill(fill_type="solid", fgColor="FFFFFF")]
        for i, r in enumerate(budget_dept, start=3):
            alloc  = r.get("allocated") or 0
            actual = r.get("actual")    or 0
            bal    = alloc - actual
            pct    = actual / alloc * 100 if alloc > 0 else 0
            for col, val in zip("ABCDE", [r.get("name") or "—", alloc, actual, bal, f"{pct:.1f}%"]):
                bws[f"{col}{i}"].value = val
                bws[f"{col}{i}"].fill  = fills[i % 2]
                bws[f"{col}{i}"].font  = Font(size=10, color=NAVY_XL)
                bws[f"{col}{i}"].alignment = Alignment(horizontal="right" if col != "A" else "left", vertical="center")
            bws.row_dimensions[i].height = 20

        # Total row
        tr = len(budget_dept) + 3
        total_alloc  = sum(r.get("allocated") or 0 for r in budget_dept)
        total_actual = sum(r.get("actual")    or 0 for r in budget_dept)
        for col, val in zip("ABCDE", ["TOTAL", total_alloc, total_actual, total_alloc - total_actual, ""]):
            bws[f"{col}{tr}"].value = val
            bws[f"{col}{tr}"].font  = Font(bold=True, size=10, color="FFFFFF")
            bws[f"{col}{tr}"].fill  = PatternFill(fill_type="solid", fgColor=NAVY_XL)
            bws[f"{col}{tr}"].alignment = Alignment(horizontal="right" if col != "A" else "left")

        # Chart % consommé
        chart_b = BarChart()
        chart_b.title = "% Budget consommé par département"
        chart_b.type  = "col"
        n = len(budget_dept)
        data_pct = Reference(bws, min_col=5, min_row=2, max_row=n + 2)
        cats_b   = Reference(bws, min_col=1, min_row=3, max_row=n + 2)
        chart_b.add_data(data_pct, titles_from_data=True)
        chart_b.set_categories(cats_b)
        chart_b.shape = 4
        bws.add_chart(chart_b, f"A{tr + 2}")

    # ── Feuille Budget Catégorie
    if budget_cat:
        currency = budget_cat[0].get("currency") or "CAD"
        cws = wb.create_sheet(title="Budget Catégorie")
        cws.column_dimensions["A"].width = 28
        for col in "BCD": cws.column_dimensions[col].width = 18

        cws.merge_cells("A1:D1")
        cws["A1"].value = f"Dépenses par catégorie — Top 10   ·   {org_name}"
        cws["A1"].font  = Font(bold=True, size=13, color="FFFFFF")
        cws["A1"].fill  = PatternFill(fill_type="solid", fgColor=NAVY_XL)
        cws["A1"].alignment = Alignment(horizontal="left", vertical="center")
        cws.row_dimensions[1].height = 28

        for col, label in zip("ABCD", ["Catégorie", f"Alloué ({currency})", f"Réel ({currency})", "% du total réel"]):
            _bhdr(cws[f"{col}2"], label, BRAND_XL)
        cws.row_dimensions[2].height = 22

        total_actual_cat = sum(r.get("actual") or 0 for r in budget_cat) or 1
        fills2 = [PatternFill(fill_type="solid", fgColor="F0F4FF"), PatternFill(fill_type="solid", fgColor="FFFFFF")]
        for i, r in enumerate(budget_cat, start=3):
            alloc  = r.get("allocated") or 0
            actual = r.get("actual")    or 0
            pct    = actual / total_actual_cat * 100
            for col, val in zip("ABCD", [r.get("name") or "—", alloc, actual, f"{pct:.1f}%"]):
                cws[f"{col}{i}"].value = val
                cws[f"{col}{i}"].fill  = fills2[i % 2]
                cws[f"{col}{i}"].font  = Font(size=10, color=NAVY_XL)
                cws[f"{col}{i}"].alignment = Alignment(horizontal="right" if col != "A" else "left", vertical="center")
            cws.row_dimensions[i].height = 20

        chart_c = PieChart()
        chart_c.title = "Répartition des dépenses par catégorie"
        nc = len(budget_cat)
        dc = Reference(cws, min_col=3, min_row=2, max_row=nc + 2)
        cc = Reference(cws, min_col=1, min_row=3, max_row=nc + 2)
        chart_c.add_data(dc, titles_from_data=True)
        chart_c.set_categories(cc)
        cws.add_chart(chart_c, f"A{nc + 4}")

    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{_filename("xlsx")}"'})


# ── Word ───────────────────────────────────────────────────────────────────

def _export_word(
    title: str,
    date_str: str,
    payload: ExportRequest,
    org_name: str,
    logo_bytes: bytes | None,
    budget_dept: list[dict],
    budget_cat: list[dict],
) -> StreamingResponse:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    NAVY_W  = RGBColor(0x0F, 0x17, 0x2A)
    BRAND_W = RGBColor(0x81, 0x8C, 0xF8)
    GRAY_W  = RGBColor(0x64, 0x74, 0x8B)

    doc = Document()
    sec = doc.sections[0]
    sec.left_margin = sec.right_margin = Cm(2.5)
    sec.top_margin  = sec.bottom_margin = Cm(2.5)

    # ── Logo + org name header
    if logo_bytes:
        try:
            p = doc.add_paragraph()
            r = p.add_run()
            r.add_picture(io.BytesIO(logo_bytes), height=Inches(0.55))
        except Exception:
            _word_org_title(doc, org_name, NAVY_W, BRAND_W)
    else:
        _word_org_title(doc, org_name, NAVY_W, BRAND_W)

    p2 = doc.add_paragraph()
    rt = p2.add_run(title); rt.font.size = Pt(18); rt.font.bold = True; rt.font.color.rgb = NAVY_W
    pd = doc.add_paragraph()
    period_word = f"Période : {payload.period_label}   ·   Généré le {date_str}" if payload.period_label else f"Généré le {date_str}"
    rd = pd.add_run(period_word); rd.font.size = Pt(10); rd.font.color.rgb = GRAY_W

    if payload.sources:
        _wlabel(doc, "SOURCES CONSULTÉES", BRAND_W)
        doc.add_paragraph(" · ".join(s.upper() for s in payload.sources))
    doc.add_paragraph()

    _wlabel(doc, "QUESTION", BRAND_W)
    doc.add_paragraph(payload.question)
    doc.add_paragraph()
    _wlabel(doc, "RÉPONSE DE L'AGENT", BRAND_W)
    for line in payload.answer.split("\n"):
        if line.strip():
            p_a = doc.add_paragraph(line); p_a.runs[0].font.size = Pt(11)
    doc.add_paragraph()

    if payload.charts:
        _wlabel(doc, "GRAPHIQUES", BRAND_W)
        for spec in payload.charts:
            try:
                img_buf = io.BytesIO(_chart_image(spec))
                doc.add_picture(img_buf, width=Inches(5.5))
                cap = doc.add_paragraph(spec.title)
                cap.runs[0].font.size = Pt(9); cap.runs[0].font.italic = True; cap.runs[0].font.color.rgb = GRAY_W
            except Exception:
                pass
        doc.add_paragraph()

    # ── Budget par département
    if budget_dept:
        currency = budget_dept[0].get("currency") or "CAD"
        total_alloc  = sum(r.get("allocated") or 0 for r in budget_dept)
        total_actual = sum(r.get("actual")    or 0 for r in budget_dept)
        _wlabel(doc, "BUDGET PAR DÉPARTEMENT — TOP 10", BRAND_W)
        doc.add_paragraph(
            f"Total alloué : {total_alloc:,.0f} {currency}   ·   "
            f"Total réel : {total_actual:,.0f} {currency}   ·   "
            f"Écart : {total_alloc - total_actual:,.0f} {currency}"
        )
        from docx.oxml.ns import qn as _qn
        tbl_d = doc.add_table(rows=1, cols=5)
        tbl_d.style = "Table Grid"
        hdr_cells = tbl_d.rows[0].cells
        for i, h in enumerate(["Département", f"Alloué ({currency})", f"Réel ({currency})", "Écart", "% Consommé"]):
            hdr_cells[i].text = h
            hdr_cells[i].paragraphs[0].runs[0].font.bold = True
            hdr_cells[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for r in budget_dept:
            alloc  = r.get("allocated") or 0
            actual = r.get("actual")    or 0
            pct    = actual / alloc * 100 if alloc > 0 else 0
            row_cells = tbl_d.add_row().cells
            for i, v in enumerate([r.get("name") or "—", f"{alloc:,.0f}", f"{actual:,.0f}", f"{alloc-actual:,.0f}", f"{pct:.1f}%"]):
                row_cells[i].text = v
        doc.add_paragraph()

    # ── Budget par catégorie
    if budget_cat:
        currency = budget_cat[0].get("currency") or "CAD"
        _wlabel(doc, "DÉPENSES PAR CATÉGORIE — TOP 10", BRAND_W)
        total_actual_cat = sum(r.get("actual") or 0 for r in budget_cat) or 1
        tbl_c = doc.add_table(rows=1, cols=4)
        tbl_c.style = "Table Grid"
        hdr2 = tbl_c.rows[0].cells
        for i, h in enumerate(["Catégorie", f"Alloué ({currency})", f"Réel ({currency})", "% du total"]):
            hdr2[i].text = h
            hdr2[i].paragraphs[0].runs[0].font.bold = True
            hdr2[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for r in budget_cat:
            actual = r.get("actual") or 0
            pct    = actual / total_actual_cat * 100
            rc = tbl_c.add_row().cells
            for i, v in enumerate([r.get("name") or "—", f"{r.get('allocated') or 0:,.0f}", f"{actual:,.0f}", f"{pct:.1f}%"]):
                rc[i].text = v
        doc.add_paragraph()

    conclusion = payload.conclusion or _default_conclusion(payload, date_str)
    _wlabel(doc, "CONCLUSION", BRAND_W)
    for line in conclusion.split("\n"):
        if line.strip():
            pc = doc.add_paragraph(line); pc.runs[0].font.size = Pt(10); pc.runs[0].font.color.rgb = GRAY_W
    doc.add_paragraph()
    pf = doc.add_paragraph()
    rf = pf.add_run(f"{org_name} · NexHire EIP · Confidentiel · {date_str}")
    rf.font.size = Pt(8); rf.font.color.rgb = GRAY_W

    buf = io.BytesIO(); doc.save(buf); buf.seek(0)
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{_filename("docx")}"'})


def _word_org_title(doc, org_name: str, navy, brand):
    p = doc.add_paragraph()
    r1 = p.add_run(org_name); r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = navy
    r2 = p.add_run("  ·  NexHire EIP"); r2.font.size = Pt(14); r2.font.bold = False; r2.font.color.rgb = brand


def _wlabel(doc, text: str, color):
    p = doc.add_paragraph()
    r = p.add_run(text); r.font.size = 9; r.font.bold = True; r.font.color.rgb = color  # type: ignore


# ── PowerPoint ─────────────────────────────────────────────────────────────

def _export_pptx(
    title: str,
    date_str: str,
    payload: ExportRequest,
    org_name: str,
    logo_bytes: bytes | None,
    budget_dept: list[dict],
    budget_cat: list[dict],
) -> StreamingResponse:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData as PptxChartData
    from pptx.enum.chart import XL_CHART_TYPE

    NAVY_P  = RGBColor(0x0F, 0x17, 0x2A)
    BRAND_P = RGBColor(0x81, 0x8C, 0xF8)
    GRAY_P  = RGBColor(0x94, 0xA3, 0xB8)
    WHITE_P = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_P = RGBColor(0xF0, 0xF4, 0xFF)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def tb(slide, l, t, w, h):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        box.text_frame.word_wrap = True
        return box

    def run(para, text, sz, bold=False, color=None, align=PP_ALIGN.LEFT):
        para.alignment = align
        r = para.add_run(); r.text = text; r.font.size = Pt(sz); r.font.bold = bold
        if color: r.font.color.rgb = color
        return r

    def add_header_bar(slide):
        bar = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(0.82))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY_P; bar.line.fill.background()
        # accent stripe
        acc = slide.shapes.add_shape(1, 0, 0, Inches(0.13), Inches(0.82))
        acc.fill.solid(); acc.fill.fore_color.rgb = BRAND_P; acc.line.fill.background()

        if logo_bytes:
            try:
                slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(0.25), Inches(0.1), height=Inches(0.62))
                return
            except Exception:
                pass
        # Text fallback
        hb = tb(slide, 0.3, 0.11, 5, 0.6)
        ph = hb.text_frame.paragraphs[0]
        run(ph, org_name[:30], 14, bold=True, color=WHITE_P)
        run(ph, " · NexHire EIP", 11, bold=False, color=BRAND_P)

    conclusion = payload.conclusion or _default_conclusion(payload, date_str)

    # ── Slide 1 : Cover
    s1 = prs.slides.add_slide(blank)
    bg = s1.background.fill; bg.solid(); bg.fore_color.rgb = NAVY_P

    if logo_bytes:
        try:
            s1.shapes.add_picture(io.BytesIO(logo_bytes), Inches(4.5), Inches(1.2), height=Inches(1.1))
        except Exception:
            bx0 = tb(s1, 1, 1.2, 11, 0.8); p0 = bx0.text_frame.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
            run(p0, org_name, 28, bold=True, color=WHITE_P)
    else:
        bx0 = tb(s1, 1, 1.2, 11, 0.8); p0 = bx0.text_frame.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        run(p0, org_name, 28, bold=True, color=WHITE_P)
        run(p0, "  ·  NexHire EIP", 18, color=BRAND_P)

    bx2 = tb(s1, 1, 3.0, 11, 1.5); p2 = bx2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    run(p2, title, 24, bold=True, color=WHITE_P)
    bx3 = tb(s1, 1, 4.8, 11, 0.5); p3 = bx3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    cover_date = f"{payload.period_label}  ·  {date_str}" if payload.period_label else date_str
    run(p3, cover_date, 13, color=GRAY_P)
    if payload.sources:
        bx4 = tb(s1, 1, 5.5, 11, 0.5); p4 = bx4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
        run(p4, " · ".join(s.upper() for s in payload.sources), 12, color=BRAND_P)

    # ── Slide 2 : Content
    s2 = prs.slides.add_slide(blank)
    bg2 = s2.background.fill; bg2.solid(); bg2.fore_color.rgb = LIGHT_P
    add_header_bar(s2)
    bt = tb(s2, 0.3, 0.95, 12.7, 0.65); p_t = bt.text_frame.paragraphs[0]
    run(p_t, title, 18, bold=True, color=NAVY_P)

    y = 1.75
    if payload.sources:
        bsr = tb(s2, 0.3, y, 12.7, 0.4); psr = bsr.text_frame.paragraphs[0]
        run(psr, "Sources : " + " · ".join(s.upper() for s in payload.sources), 10, color=BRAND_P)
        y += 0.5

    bq = tb(s2, 0.3, y, 12.7, 0.5); pq = bq.text_frame.paragraphs[0]
    run(pq, f"Q : {payload.question[:200]}", 10, color=GRAY_P); y += 0.6

    ba = tb(s2, 0.3, y, 12.7, 7.5 - y - 0.2); tf = ba.text_frame
    for i, line in enumerate([l for l in payload.answer.split("\n") if l.strip()][:20]):
        p_ans = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p_ans.space_before = Pt(1)
        r = p_ans.add_run(); r.text = line; r.font.size = Pt(10.5); r.font.color.rgb = NAVY_P

    # ── Slides 3+ : Charts
    for spec in payload.charts:
        try:
            sg = prs.slides.add_slide(blank)
            bgg = sg.background.fill; bgg.solid(); bgg.fore_color.rgb = LIGHT_P
            add_header_bar(sg)
            bct = tb(sg, 0.3, 0.95, 12.7, 0.65); p_ct = bct.text_frame.paragraphs[0]
            run(p_ct, spec.title, 18, bold=True, color=NAVY_P)

            cd = PptxChartData()
            cd.categories = spec.labels
            cd.add_series(spec.title, [float(v) for v in spec.values])
            ct = (XL_CHART_TYPE.PIE if spec.type == "pie"
                  else XL_CHART_TYPE.LINE if spec.type == "line"
                  else XL_CHART_TYPE.COLUMN_CLUSTERED)
            chart = sg.shapes.add_chart(ct, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5), cd)
            ch = chart.chart
            ch.has_legend = spec.type == "pie"
            if ch.has_legend: ch.legend.position = 2
        except Exception:
            pass

    # ── Slide : Budget par département
    if budget_dept:
        try:
            currency = budget_dept[0].get("currency") or "CAD"
            sd = prs.slides.add_slide(blank)
            bgd = sd.background.fill; bgd.solid(); bgd.fore_color.rgb = LIGHT_P
            add_header_bar(sd)
            btd = tb(sd, 0.3, 0.95, 12.7, 0.65); p_td = btd.text_frame.paragraphs[0]
            run(p_td, "Budget par département — Top 10", 18, bold=True, color=NAVY_P)

            cd_b = PptxChartData()
            cd_b.categories = [r.get("name") or "—" for r in budget_dept]
            pct_vals = [(r.get("actual") or 0) / (r.get("allocated") or 1) * 100 for r in budget_dept]
            cd_b.add_series("% Consommé", [round(v, 1) for v in pct_vals])
            ch_b = sd.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5), cd_b)
            ch_b.chart.has_legend = False
        except Exception:
            pass

    # ── Slide : Dépenses par catégorie
    if budget_cat:
        try:
            sc2 = prs.slides.add_slide(blank)
            bgc2 = sc2.background.fill; bgc2.solid(); bgc2.fore_color.rgb = LIGHT_P
            add_header_bar(sc2)
            btc2 = tb(sc2, 0.3, 0.95, 12.7, 0.65); p_tc2 = btc2.text_frame.paragraphs[0]
            run(p_tc2, "Dépenses par catégorie — Top 10", 18, bold=True, color=NAVY_P)

            cd_c = PptxChartData()
            cd_c.categories = [r.get("name") or "—" for r in budget_cat]
            cd_c.add_series("Dépenses réelles", [float(r.get("actual") or 0) for r in budget_cat])
            ch_c = sc2.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(1.5), Inches(1.5), Inches(10), Inches(5.5), cd_c)
            ch_c.chart.has_legend = True
        except Exception:
            pass

    # ── Slide : Conclusion
    sc = prs.slides.add_slide(blank)
    bgc = sc.background.fill; bgc.solid(); bgc.fore_color.rgb = NAVY_P
    add_header_bar(sc)
    bct2 = tb(sc, 0.5, 1.0, 12.3, 0.7); pct2 = bct2.text_frame.paragraphs[0]
    run(pct2, "Conclusion", 22, bold=True, color=WHITE_P)
    bcb = tb(sc, 0.5, 2.0, 12.3, 4.5); tf_c = bcb.text_frame
    for i, line in enumerate([l for l in conclusion.split("\n") if l.strip()]):
        p_c = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p_c.space_before = Pt(4); r = p_c.add_run()
        r.text = line; r.font.size = Pt(12); r.font.color.rgb = GRAY_P
    bdate = tb(sc, 0.5, 6.8, 12.3, 0.4); pd = bdate.text_frame.paragraphs[0]; pd.alignment = PP_ALIGN.CENTER
    run(pd, f"{org_name} · NexHire EIP · Confidentiel · {date_str}", 9, color=GRAY_P)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{_filename("pptx")}"'})
